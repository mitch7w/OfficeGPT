from flask import Flask, request, send_file
from flask_cors import CORS
from docx import *
from openai import OpenAI
from dotenv import load_dotenv
import io
import os
import xlsxwriter
from pptx import *

# setup
load_dotenv()
OPENAI_API_KEY = os.environ.get('OPENAI_API_KEY')
app = Flask(__name__)
CORS(app)  # Enable CORS for all routes
client = OpenAI()


def openai_call(document_instructions, system_instructions):
    completion = client.chat.completions.create(
        model="gpt-4-1106-preview",
        messages=[
            {"role": "system", "content": system_instructions},
            {"role": "user", "content": document_instructions}
        ]
    )
    return completion.choices[0].message.content


def clean_commands(gpt_output):
    # prompt injection checking for malicious python commands
    print("gpt_output: ", gpt_output)
    completion = client.chat.completions.create(model="gpt-3.5-turbo-1106", messages=[
        {"role": "system",
            "content": "You are an expert Python code reviewer. Please check the following Python code for any malicious Python commands that could affect a user's file system, expose system or user data or in any way harm a user or their computer. Please remove any malicious code you find. Additionally please look for any missing library imports and insert them at the beginning of the code. Finally - please modify any slide.placeholders[1] to slide.placeholders[0]. Please return the python code and only the python code - nothing else."},
        {"role": "user", "content": gpt_output}
    ]
    )
    secure_output = completion.choices[0].message.content
    print("secure_output: ", secure_output)
    lines = secure_output.splitlines()
    code_lines = []
    for line in lines:
        if line.startswith("```python") or line.startswith("```"):
            continue
        else:
            code_lines.append(line)
    cleaned_code = '\n'.join(code_lines)
    return cleaned_code


# creates a docx file based on the OpenAI instructions
def create_docx(document_instructions):
    doc_buffer = io.BytesIO()  # Save the created document to a BytesIO buffer
    document = Document()
    system_instructions = "You are a Word document expert and create word documents using only commands from the python-docx library. Please do not return any text - only the Python code. I have already ran the command from docx import * and setup the document with document = Document() as well as document.save() so do not include these commands in your response. Please however import any extra libraries, constants or additional components of the library as needed. So I do not need any functions or other code or any extra ``` or designations that this is python code - return only the python-docx commands for creating the specific items in user's document. Please think step-by-step, use only zero-based indexing for arrays - making sure items exist before you try referencing them and only include working Python code that will compile, does not contain undefined variables, is complete and is error-free. Thank you."
    # get python-docx commands from GPT
    creation_commands = openai_call(document_instructions, system_instructions)

    cleaned_commands = clean_commands(creation_commands)
    print(cleaned_commands, flush=True)
    # actually run the commands that create the Word doc elements
    exec(cleaned_commands)

    document.save(doc_buffer)  # save doc to this buffer variable
    doc_buffer.seek(0)
    return doc_buffer


def create_excel(document_instructions):
    excel_buffer = io.BytesIO()  # Save the Excel workbook to a BytesIO buffer
    workbook = xlsxwriter.Workbook(excel_buffer)
    system_instructions = "You are an Excel document expert and create Excel documents using only commands from the XlsxWriter library. Please do not return any text - only the Python code. I have already ran the command import xlsxwriter and setup the document with worksheet = workbook.add_worksheet() as well as workbook.close() so do not include these commands in your response. Please however import any extra libraries, constants or additional components of the library as needed. So I do not need any functions or other code or any extra ``` or designations that this is python code - return only the XlsxWriter commands for creating the specific items in user's Excel document. Please think step-by-step, use only zero-based indexing for arrays - making sure items exist before you try referencing them and only include working Python code that will compile, does not contain undefined variables, is complete and is error-free. Thank you."
    worksheet = workbook.add_worksheet()
    # get XlsxWriter commands from GPT
    creation_commands = openai_call(document_instructions, system_instructions)

    cleaned_commands = clean_commands(creation_commands)
    print(cleaned_commands, flush=True)
    # actually run the commands that create the Excel doc elements
    exec(cleaned_commands)

    workbook.close()
    excel_buffer.seek(0)
    return excel_buffer


def create_powerpoint(document_instructions):
    ppt_buffer = io.BytesIO()  # Create an in-memory buffer
    prs = Presentation()
    system_instructions = "You are a PowerPoint document expert and create PowerPoint documents using only commands from the python-pptx library. Please do not return any text - only the Python code. I have already setup the document with prs = Presentation() as well as prs.save(ppt_buffer) so do not include these commands in your response. Please however import any extra libraries, constants or additional components of the library as needed. So I do not need any functions or other code or any extra ``` or designations that this is python code - return only the python-pptx commands for creating the specific items in user's PowerPoint document. Please do not insert any images. Please think step-by-step, use only zero-based indexing for arrays - making sure items exist before you try referencing them and only include working Python code that will compile, does not contain undefined variables, is complete and is error-free. Thank you."
    # get python-pptx commands from GPT
    creation_commands = openai_call(document_instructions, system_instructions)

    cleaned_commands = clean_commands(creation_commands)
    print(cleaned_commands, flush=True)
    # actually run the commands that create the PowerPoint doc elements
    exec(cleaned_commands)

    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer


@app.route('/create', methods=['GET', 'POST'])
def create_endpoint():
    document_instructions = request.args.get(
        'document_instructions', default='', type=str)
    selected_product = request.args.get(
        'selected_product', default='', type=str)

    download_name = None
    mimetype = None
    doc_buffer = None

    # create the document using the correct function and send it back to the front-end
    if (selected_product == "Word"):
        # Generate the Word document
        doc_buffer = create_docx(document_instructions)
        download_name = 'created.docx'
        mimetype = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
    elif (selected_product == "Excel"):
        # Generate the Word document
        doc_buffer = create_excel(document_instructions)
        download_name = 'created.xlsx'
        mimetype = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    elif (selected_product == "PowerPoint"):
        # Generate the Word document
        doc_buffer = create_powerpoint(document_instructions)
        download_name = 'created.pptx'
        mimetype = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'

    return send_file(
        doc_buffer,
        as_attachment=True,
        download_name=download_name,
        mimetype=mimetype)
