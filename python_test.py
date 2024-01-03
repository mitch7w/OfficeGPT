from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_LEGEND_POSITION, XL_CHART_TYPE
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# create Presentation object
prs = Presentation()

# Add a title slide with the main topic
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[0]
title.text = "National Parks of the USA"
subtitle.text = "An overview of prominent national parks and their statistics"

# Add a slide with a list of National Parks
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "List of National Parks"
content = slide.placeholders[0]
content.text = "Yosemite\nYellowstone\nGrand Canyon\nZion\nGreat Smoky Mountains"

# Add a slide with a table of statistics
slide_layout = prs.slide_layouts[5]  # Title Only layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Statistics of National Parks"

rows, cols = 6, 4
left = Inches(2)
top = Inches(2)
width = Inches(6)
height = Inches(0.8)
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Set column widths
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(1.5)
table.columns[2].width = Inches(1.5)
table.columns[3].width = Inches(1.0)

# Write header row
table.cell(0, 0).text = 'National Park'
table.cell(0, 1).text = 'Visitors'
table.cell(0, 2).text = 'Area'
table.cell(0, 3).text = 'Established'
prs.save("hello.pptx")
